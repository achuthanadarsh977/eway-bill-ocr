from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy
ACCENT     = RGBColor(0x5C, 0x9E, 0xFF)   # blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
ORANGE     = RGBColor(0xFF, 0x98, 0x00)
CARD_BG    = RGBColor(0x2A, 0x2A, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper functions ──────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, color=CARD_BG, radius=False):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def divider(slide, top, color=ACCENT):
    """Thin horizontal rule."""
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(top), Inches(12.33), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def bullet_box(slide, items, left, top, width, height, size=16, color=LIGHT_GREY, indent="  •  "):
    """Multi-line bullet list in a text box."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{indent}{item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.6, ACCENT)
txt(s, "OCR MODULE DESIGN", 0.5, 0.12, 12, 0.5,
    size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(s, "Eway Bill Document Parser", 0.5, 1.2, 12, 1.0,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(s, 2.8)

txt(s, "Converts Eway Bill PDFs → Structured JSON automatically",
    0.5, 3.0, 12.33, 0.6, size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

for i, (label, val) in enumerate([
    ("Prepared by", "Adarsh"),
    ("Date", "04-Jun-2026"),
    ("Status", "Draft — Pending Review"),
]):
    x = 1.5 + i * 3.8
    box(s, x, 4.2, 3.2, 1.2, CARD_BG)
    txt(s, label, x + 0.1, 4.3, 3.0, 0.35, size=13, color=ACCENT, bold=True)
    txt(s, val,   x + 0.1, 4.7, 3.0, 0.55, size=16, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Objective
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "01  |  OBJECTIVE", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

txt(s, "What are we building?", 0.5, 0.9, 12, 0.55,
    size=28, bold=True, color=WHITE)
divider(s, 1.6)

goals = [
    "Accept an Eway Bill PDF as input",
    "Detect if PDF is digital or scanned automatically",
    "Extract all relevant fields using OCR",
    "Return structured data in JSON format",
    "Plug into any downstream application — fully independent module",
]
for i, goal in enumerate(goals):
    y = 1.8 + i * 0.9
    box(s, 0.5, y, 0.55, 0.55, ACCENT)
    txt(s, str(i + 1), 0.5, y, 0.55, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, goal, 1.2, y + 0.05, 11.5, 0.5, size=18, color=WHITE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Flow
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "02  |  HIGH-LEVEL DATA FLOW", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

stages = [
    ("Input\nPDF",     "📄"),
    ("PDF\nLoader",    "🔍"),
    ("OCR\nEngine",    "⚙️"),
    ("Parser",         "📋"),
    ("Formatter",      "📤"),
]
labels = ["", "Detect type", "Extract text", "Pull fields", "JSON / CSV"]
colors = [RGBColor(0x37, 0x47, 0x6F), ACCENT, RGBColor(0x7E, 0x57, 0xC2),
          RGBColor(0x26, 0xA6, 0x9A), GREEN]

for i, ((name, icon), label, col) in enumerate(zip(stages, labels, colors)):
    x = 0.4 + i * 2.5
    box(s, x, 2.0, 2.1, 1.8, col)
    txt(s, icon, x, 2.1, 2.1, 0.6, size=28, align=PP_ALIGN.CENTER)
    txt(s, name, x, 2.7, 2.1, 0.7, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if label:
        txt(s, label, x, 3.9, 2.1, 0.4, size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)
    if i < 4:
        txt(s, "▶", x + 2.1, 2.6, 0.4, 0.6, size=22, color=ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Module 1 & 2
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "03  |  MODULES 1 & 2", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

# Module 1
box(s, 0.3, 0.9, 6.0, 5.8, CARD_BG)
txt(s, "Module 1 — PDF Loader", 0.5, 1.0, 5.6, 0.5, size=20, bold=True, color=ACCENT)
divider(s, 1.65)
bullet_box(s, [
    "Load the PDF from file path",
    "Attempt text extraction with pdf-parse",
    "Text > 50 chars  →  Digital PDF",
    "Text empty/minimal  →  Scanned PDF",
    "Returns: path, numPages, isDigital,",
    "rawBuffer, pagesText[]",
], 0.4, 1.75, 5.8, 4.5)

# Module 2
box(s, 7.0, 0.9, 6.0, 5.8, CARD_BG)
txt(s, "Module 2 — OCR Engine", 7.2, 1.0, 5.6, 0.5, size=20, bold=True, color=ACCENT)
divider(s, 1.65)
bullet_box(s, [
    "Routes based on PDF type:",
    "",
    "Digital PDF:",
    "   → Return text directly (fast)",
    "",
    "Scanned PDF:",
    "   → Convert pages to images",
    "   → Run tesseract.js (WASM)",
    "   → Return extracted text",
], 7.1, 1.75, 5.8, 4.5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Module 3 (Parser Fields)
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "04  |  MODULE 3 — PARSER", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

txt(s, "Fields extracted from every Eway Bill", 0.5, 0.85, 12, 0.5,
    size=22, bold=True, color=WHITE)
divider(s, 1.5)

sections = [
    ("EWB Details",        ["EWB Number", "EWB Date", "Valid Until", "Transaction Type",
                             "Document No & Date", "Value of Goods", "IRN"]),
    ("From Party",         ["GSTIN", "Company Name", "Place", "State", "Pincode"]),
    ("To Party",           ["GSTIN", "Company Name", "Place", "State", "Pincode"]),
    ("Item Details",       ["HSN Code", "Product Description"]),
    ("Transporter",        ["Transporter GSTIN", "Transporter Name",
                             "Vehicle Number", "Transport Mode"]),
]
colors2 = [ACCENT, RGBColor(0x26, 0xA6, 0x9A), RGBColor(0x7E, 0x57, 0xC2),
           ORANGE, GREEN]

col_x = [0.3, 2.95, 5.6, 8.25, 10.9]
for i, ((sec, fields), col, cx) in enumerate(zip(sections, colors2, col_x)):
    box(s, cx, 1.65, 2.3, 5.3, CARD_BG)
    box(s, cx, 1.65, 2.3, 0.55, col)
    txt(s, sec, cx + 0.1, 1.7, 2.1, 0.45, size=13, bold=True, color=WHITE)
    bullet_box(s, fields, cx + 0.1, 2.3, 2.1, 4.5, size=13, indent="• ")


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Module 4 & 5
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "05  |  MODULES 4 & 5", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

# Module 4
box(s, 0.3, 0.9, 6.0, 5.8, CARD_BG)
txt(s, "Module 4 — Formatter", 0.5, 1.0, 5.6, 0.5, size=20, bold=True, color=ACCENT)
divider(s, 1.65)
bullet_box(s, [
    "Saves parsed data to disk",
    "",
    "JSON  (primary)",
    "   → Pretty-printed nested structure",
    "   → Used by downstream apps",
    "",
    "CSV  (secondary)",
    "   → Flattened single-row format",
    "   → For manual review / Excel",
    "",
    "Output: {filename}_{timestamp}.json",
], 0.4, 1.75, 5.8, 4.5)

# Module 5
box(s, 7.0, 0.9, 6.0, 5.8, CARD_BG)
txt(s, "Module 5 — Main Entry Point", 7.2, 1.0, 5.6, 0.5, size=20, bold=True, color=ACCENT)
divider(s, 1.65)
bullet_box(s, [
    "Wires all 4 modules together",
    "Accepts CLI arguments",
    "",
    "Commands:",
    "  ts-node main.ts <pdf>",
    "     → outputs JSON",
    "  ts-node main.ts <pdf> --csv",
    "     → also outputs CSV",
    "  ts-node main.ts <pdf> --raw",
    "     → prints raw text",
], 7.1, 1.75, 5.8, 4.5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "06  |  TECHNOLOGY STACK", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

stack = [
    ("TypeScript",    "Language",                   "Type safety, better maintainability"),
    ("Node.js",       "Runtime",                    "Cross-platform, no server needed"),
    ("pdf-parse",     "Digital PDF Extraction",     "Pure JS — no system dependencies"),
    ("tesseract.js",  "Scanned PDF OCR",            "WASM-based — no system install"),
    ("JSON + CSV",    "Output Formats",             "JSON for apps, CSV for manual review"),
]
RH   = 0.95   # row height — taller so text has breathing room
STEP = 1.08   # step between rows
for i, (tech, role, reason) in enumerate(stack):
    y = 0.9 + i * STEP
    box(s, 0.3,  y, 2.5, RH, ACCENT)
    txt(s, tech, 0.4, y + 0.25, 2.3, 0.55, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 2.85, y, 4.0, RH, CARD_BG)
    txt(s, role, 3.0, y + 0.28, 3.7, 0.5, size=16, color=LIGHT_GREY)
    box(s, 6.9,  y, 6.1, RH, CARD_BG)
    txt(s, reason, 7.05, y + 0.28, 5.8, 0.5, size=15, color=WHITE)

txt(s, "Why no cloud API?  →  Free, local, no rate limits, sustainable at any volume",
    0.5, 6.6, 12.33, 0.45, size=15, color=ACCENT, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Build Plan
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "07  |  INCREMENTAL BUILD PLAN", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

txt(s, "Phases 1–3 complete  •  Phases 4–6 pending",
    0.3, 0.85, 12.73, 0.45, size=15, color=LIGHT_GREY, italic=True)
divider(s, 1.45)

phases = [
    ("Phase 1", "PDF Loader + Digital text extraction",    "DONE",    GREEN),
    ("Phase 2", "Parser — all Eway Bill fields via regex", "DONE",    GREEN),
    ("Phase 3", "JSON + CSV formatter",                   "DONE",    GREEN),
    ("Phase 4", "Scanned PDF support via tesseract.js",   "PENDING", ORANGE),
    ("Phase 5", "Batch processing — folder of PDFs",      "PENDING", ORANGE),
    ("Phase 6", "REST API wrapper (optional)",            "PENDING", ORANGE),
]
PRH  = 0.75   # phase row height
PSTEP = 0.88  # step between rows — 0.13" gap between rows
for i, (phase, desc, status, col) in enumerate(phases):
    y = 1.6 + i * PSTEP
    box(s, 0.3,   y, 1.6,  PRH, col)
    txt(s, phase, 0.35, y + 0.22, 1.5, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 2.0,   y, 8.5,  PRH, CARD_BG)
    txt(s, desc,  2.15, y + 0.22, 8.2, 0.4, size=15, color=WHITE)
    box(s, 10.6,  y, 2.4,  PRH, col)
    txt(s, status, 10.65, y + 0.22, 2.3, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Open Questions
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.7, ACCENT)
txt(s, "08  |  OPEN QUESTIONS FOR REVIEW", 0.5, 0.15, 12, 0.5, size=20, bold=True, color=WHITE)

txt(s, "Seeking answers before Phase 4 implementation begins",
    0.5, 0.85, 12, 0.42, size=16, color=LIGHT_GREY, italic=True)
divider(s, 1.42)

questions = [
    ("Q1", "Multi-item Eway Bills",
     "Should the module handle documents with multiple HSN codes / items per bill?"),
    ("Q2", "API vs CLI",
     "Do we need a REST API wrapper, or is CLI usage sufficient for now?"),
    ("Q3", "JSON Schema",
     "Should the output JSON follow a specific schema already used in the system?"),
    ("Q4", "Future Document Types",
     "Are there more documents to handle — invoices, delivery challans, etc.?"),
]
QH   = 1.1    # question box height — taller for breathing room
QSTEP = 1.28  # step — 0.18" gap between question rows
for i, (q, title, detail) in enumerate(questions):
    y = 1.6 + i * QSTEP
    box(s, 0.3,  y, 0.85, QH, ACCENT)
    txt(s, q,    0.3,  y + 0.32, 0.85, 0.55, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.25, y, 11.78, QH, CARD_BG)
    txt(s, title,  1.4, y + 0.12, 11.3, 0.38, size=16, bold=True, color=ACCENT)
    txt(s, detail, 1.4, y + 0.58, 11.3, 0.42, size=14, color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You / Summary
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

box(s, 0, 0, 13.33, 0.6, ACCENT)

txt(s, "Ready for Review", 0.5, 1.0, 12.33, 1.1,
    size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, 2.45)

txt(s, "Phases 1–3 are complete and working on sample Eway Bill PDFs",
    0.5, 2.62, 12.33, 0.55, size=19, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

summary = [
    ("5 Modules",        "Independent, plug-and-play"),
    ("2 Output Formats", "JSON + CSV"),
    ("0 Cloud APIs",     "Fully local, sustainable"),
    ("3 Phases Done",    "Working on real PDFs today"),
]
CARD_Y  = 3.4    # closer to subtitle — was 3.7
CARD_H  = 2.0    # taller cards — was 1.8
for i, (stat, label) in enumerate(summary):
    x = 0.9 + i * 3.1
    box(s, x, CARD_Y, 2.8, CARD_H, CARD_BG)
    # accent strip at top of card
    box(s, x, CARD_Y, 2.8, 0.12, ACCENT)
    # stat text centred in upper half of card
    txt(s, stat,  x + 0.1, CARD_Y + 0.55, 2.6, 0.65,
        size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # label centred in lower half
    txt(s, label, x + 0.1, CARD_Y + 1.3,  2.6, 0.55,
        size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

txt(s, "Prepared by Adarsh  |  04-Jun-2026  |  Pending Manager Review",
    0.5, 6.9, 12.33, 0.42, size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)


# ── Save ──────────────────────────────────────────────────────────
out = r"C:\Users\SriniAchuthan\Documents\eway_bill_ocr_ts\Eway_Bill_OCR_Design.pptx"
prs.save(out)
print(f"Saved: {out}")
